"""
PPTX Conversion Node — Converts HTML to editable .pptx file.
Calls the existing convert.py script or uses python-pptx directly.

Model strategy: No LLM needed — deterministic conversion.
"""

import subprocess
import sys
from pathlib import Path


def convert_html_to_pptx(html_path: str, theme: str = "dark") -> str:
    """
    Convert HTML presentation to .pptx using the existing convert.py script.
    Returns path to generated .pptx file.
    """
    html_path = Path(html_path)
    if not html_path.exists():
        raise FileNotFoundError(f"HTML file not found: {html_path}")

    pptx_path = html_path.with_suffix(".pptx")

    # Determine theme from template name
    pptx_theme = "dark"
    if "light" in theme.lower() or "corporate" in theme.lower():
        pptx_theme = "light"

    # Call the convert.py script
    convert_script = Path(".agent/skills/html-to-pptx/scripts/convert.py")

    if convert_script.exists():
        result = subprocess.run(
            [sys.executable, str(convert_script), str(html_path), str(pptx_path), pptx_theme],
            capture_output=True,
            text=True,
            timeout=60,
        )

        if result.returncode == 0 and pptx_path.exists():
            return str(pptx_path)
        else:
            # Fallback: generate directly
            return generate_pptx_direct(html_path, pptx_path, pptx_theme)
    else:
        return generate_pptx_direct(html_path, pptx_path, pptx_theme)


def generate_pptx_direct(html_path: Path, pptx_path: Path, theme: str = "dark") -> str:
    """
    Direct PPTX generation from HTML using python-pptx.
    Fallback if convert.py is not available.
    """
    try:
        # Import the convert module directly
        sys.path.insert(0, str(Path(".agent/skills/html-to-pptx/scripts")))
        from convert import convert

        result = convert(str(html_path), str(pptx_path), theme)
        return result
    except ImportError:
        # Last resort: minimal PPTX from state
        return _minimal_pptx_from_html(html_path, pptx_path, theme)


def _minimal_pptx_from_html(html_path: Path, pptx_path: Path, theme: str) -> str:
    """Minimal PPTX generation as last resort fallback."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from bs4 import BeautifulSoup

    with open(html_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "lxml")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    if theme == "dark":
        bg_color = RGBColor(0x0F, 0x17, 0x2A)
        title_color = RGBColor(0xF1, 0xF5, 0xF9)
        text_color = RGBColor(0xCB, 0xD5, 0xE1)
    else:
        bg_color = RGBColor(0xFF, 0xFF, 0xFF)
        title_color = RGBColor(0x1E, 0x29, 0x3B)
        text_color = RGBColor(0x33, 0x41, 0x55)

    # Find slides
    slide_elements = soup.select(".slide, [id^='slide-']")

    for slide_el in slide_elements:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        y_pos = Inches(0.5)

        # Title
        title_el = slide_el.find(["h1", "h2"])
        if title_el:
            txBox = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11.5), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_el.get_text(strip=True)
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = title_color
            y_pos += Inches(1.0)

        # Bullets
        bullets = slide_el.find_all("li")
        if bullets:
            txBox = slide.shapes.add_textbox(
                Inches(0.8), y_pos, Inches(11.5), Inches(len(bullets) * 0.4)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, li in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {li.get_text(strip=True)}"
                p.font.size = Pt(13)
                p.font.color.rgb = text_color

    prs.save(str(pptx_path))
    return str(pptx_path)


def pptx_node(state: dict) -> dict:
    """
    LangGraph node: PPTX Conversion.
    Converts the generated HTML into an editable .pptx file.

    No LLM needed — deterministic conversion.

    Input: html_path, template
    Output: pptx_path
    """
    html_path = state.get("html_path", "")
    template = state.get("template", "Padrão (dark theme)")
    formato_saida = state.get("formato_saida", [])

    # Skip if PPTX not requested
    if "PPTX (editável)" not in formato_saida:
        return {**state, "pptx_path": ""}

    if not html_path or not Path(html_path).exists():
        return {**state, "pptx_path": "", "error": "HTML não encontrado para conversão"}

    try:
        pptx_path = convert_html_to_pptx(html_path, template)

        if pptx_path and Path(pptx_path).exists():
            return {**state, "pptx_path": pptx_path}
        else:
            return {**state, "pptx_path": "", "error": "Falha na conversão PPTX"}

    except Exception as e:
        return {**state, "pptx_path": "", "error": f"Erro na conversão: {str(e)}"}
