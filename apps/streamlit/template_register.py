"""
Corporate Template System — Template registration via .pptx upload.
Deterministic extraction of colors, fonts, and logo from slide master.

No LLM needed — pure python-pptx metadata extraction.
"""

import json
import shutil
import re
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


def extract_theme_colors(prs: Presentation) -> dict:
    """
    Extract color palette from the PPTX theme/slide master.
    Returns dict with primary, secondary, accent, background, text colors.
    """
    colors = {
        "primary": "#0066CC",
        "secondary": "#003366",
        "accent": "#00AAFF",
        "background": "#FFFFFF",
        "text": "#1A1A1A",
    }

    try:
        # Try to get colors from theme
        slide_master = prs.slide_masters[0]
        theme = slide_master.slide_layouts[0]

        # Extract from slide master background
        bg = slide_master.background
        if bg.fill.type is not None:
            try:
                fg_color = bg.fill.fore_color
                if fg_color.type is not None:
                    colors["background"] = f"#{fg_color.rgb}" if fg_color.rgb else colors["background"]
            except Exception:
                pass

        # Try to extract from shapes in the master
        for shape in slide_master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            rgb = str(run.font.color.rgb)
                            if rgb != "000000":
                                colors["text"] = f"#{rgb}"
                                break

    except Exception:
        pass

    # Try to extract accent colors from first slide content
    try:
        if prs.slides:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.rgb:
                                rgb = str(run.font.color.rgb)
                                hex_color = f"#{rgb}"
                                # Heuristic: bright colors are likely accent/primary
                                if rgb not in ("000000", "FFFFFF", "333333"):
                                    if colors["primary"] == "#0066CC":
                                        colors["primary"] = hex_color
                                    elif colors["accent"] == "#00AAFF":
                                        colors["accent"] = hex_color
    except Exception:
        pass

    return colors


def extract_fonts(prs: Presentation) -> dict:
    """Extract font families from the PPTX."""
    fonts = {"title": "Arial", "body": "Arial"}

    try:
        # Check slide master for fonts
        slide_master = prs.slide_masters[0]
        for shape in slide_master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            # Larger fonts are likely titles
                            if run.font.size and run.font.size >= Pt(20):
                                fonts["title"] = run.font.name
                            else:
                                fonts["body"] = run.font.name

        # Also check first slide
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and run.font.size:
                                if run.font.size >= Pt(24):
                                    fonts["title"] = run.font.name
                                elif fonts["body"] == "Arial":
                                    fonts["body"] = run.font.name

    except Exception:
        pass

    return fonts


def extract_logo(prs: Presentation, output_dir: Path) -> Optional[str]:
    """
    Extract logo image from slide master or first slide.
    Saves to output_dir and returns filename.
    """
    try:
        # Check slide master for images (logos are usually there)
        slide_master = prs.slide_masters[0]
        for shape in slide_master.shapes:
            if shape.shape_type == 13:  # Picture
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                logo_filename = f"logo.{ext}"
                logo_path = output_dir / logo_filename
                with open(logo_path, "wb") as f:
                    f.write(image.blob)
                return logo_filename

        # Fallback: check first slide for small images (likely logo)
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if shape.shape_type == 13:  # Picture
                    # Heuristic: logos are small (< 3 inches wide)
                    if shape.width and shape.width < 3000000:  # ~3 inches in EMU
                        image = shape.image
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        logo_filename = f"logo.{ext}"
                        logo_path = output_dir / logo_filename
                        with open(logo_path, "wb") as f:
                            f.write(image.blob)
                        return logo_filename

    except Exception:
        pass

    return None


def get_slide_layouts(prs: Presentation) -> list[str]:
    """List available slide layout names."""
    layouts = []
    try:
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                name = layout.name
                if name and name not in layouts:
                    layouts.append(name)
    except Exception:
        pass
    return layouts or ["Blank"]


def register_template(
    pptx_file_path: str,
    template_name: str,
    template_id: Optional[str] = None,
    description: str = "",
    footer: str = "",
) -> dict:
    """
    Register a new corporate template from a .pptx file.

    1. Extracts colors, fonts, logo from the PPTX
    2. Copies the .pptx as base.pptx
    3. Generates manifest.json
    4. Updates registry.json

    Returns the generated manifest dict.
    """
    from config import get_settings

    settings = get_settings()
    templates_dir = settings.templates_dir

    # Generate ID from name
    if not template_id:
        template_id = re.sub(r"[^\w\s-]", "", template_name.lower())
        template_id = re.sub(r"[\s]+", "-", template_id)

    # Create template directory
    template_dir = templates_dir / template_id
    template_dir.mkdir(parents=True, exist_ok=True)

    # Load PPTX
    prs = Presentation(pptx_file_path)

    # Extract metadata
    colors = extract_theme_colors(prs)
    fonts = extract_fonts(prs)
    logo = extract_logo(prs, template_dir)
    layouts = get_slide_layouts(prs)

    # Copy base.pptx
    base_path = template_dir / "base.pptx"
    shutil.copy2(pptx_file_path, base_path)

    # Generate manifest
    manifest = {
        "id": template_id,
        "name": template_name,
        "description": description or f"Template corporativo: {template_name}",
        "colors": colors,
        "fonts": fonts,
        "logo": logo,
        "footer": footer,
        "layouts": layouts,
    }

    # Save manifest
    manifest_path = template_dir / "manifest.json"
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)

    # Update registry
    registry_path = templates_dir / "registry.json"
    registry = []
    if registry_path.exists():
        with open(registry_path, "r", encoding="utf-8") as f:
            registry = json.load(f)

    # Remove existing entry with same ID
    registry = [t for t in registry if t.get("id") != template_id]

    # Add new entry
    registry.append({
        "id": template_id,
        "name": template_name,
        "description": manifest["description"],
        "colors": colors,
        "fonts": fonts,
    })

    with open(registry_path, "w", encoding="utf-8") as f:
        json.dump(registry, f, indent=2, ensure_ascii=False)

    return manifest


def list_templates() -> list[dict]:
    """List all registered templates."""
    from config import get_settings

    settings = get_settings()
    registry_path = settings.templates_dir / "registry.json"

    if registry_path.exists():
        with open(registry_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return []


def delete_template(template_id: str) -> bool:
    """Remove a template from registry and delete its directory."""
    from config import get_settings

    settings = get_settings()
    templates_dir = settings.templates_dir

    # Remove directory
    template_dir = templates_dir / template_id
    if template_dir.exists():
        shutil.rmtree(template_dir)

    # Update registry
    registry_path = templates_dir / "registry.json"
    if registry_path.exists():
        with open(registry_path, "r", encoding="utf-8") as f:
            registry = json.load(f)
        registry = [t for t in registry if t.get("id") != template_id]
        with open(registry_path, "w", encoding="utf-8") as f:
            json.dump(registry, f, indent=2, ensure_ascii=False)

    return True
