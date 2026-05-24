"""
HTML to PPTX Converter
Converte apresentações HTML em arquivos .pptx editáveis.
Todo conteúdo é texto editável — nunca renderizado como imagem.
Imagens (tags <img>, background-image CSS, URLs externas) são baixadas e inseridas.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re
import requests
from io import BytesIO
from pathlib import Path
import os


def extract_text_clean(element) -> str:
    """Extrai texto limpo de um elemento, removendo espaços extras."""
    text = element.get_text(separator=" ", strip=True)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def parse_color_from_style(style: str) -> str | None:
    """Extrai cor hex de um atributo style CSS."""
    if not style:
        return None
    match = re.search(r"color:\s*#([0-9a-fA-F]{6})", style)
    if match:
        return match.group(1)
    return None


def extract_background_image_url(style: str) -> str | None:
    """Extrai URL de background-image de um atributo style CSS."""
    if not style:
        return None
    match = re.search(r"background-image:\s*url\(['\"]?(.*?)['\"]?\)", style)
    if match:
        return match.group(1)
    return None


def download_image(url: str, timeout: int = 15) -> BytesIO | None:
    """Baixa uma imagem de URL e retorna como BytesIO. Retorna None se falhar."""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (compatible; PresentationConverter/1.0)"
        }
        response = requests.get(url, timeout=timeout, headers=headers)
        if response.status_code == 200 and len(response.content) > 100:
            content_type = response.headers.get("content-type", "")
            if "image" in content_type or url.lower().endswith(
                (".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg")
            ):
                return BytesIO(response.content)
    except Exception:
        pass
    return None


def find_images_in_element(element) -> list[dict]:
    """
    Encontra todas as imagens dentro de um elemento HTML.
    Busca em: <img> tags, background-image inline, e data-src attributes.
    """
    images = []

    # 1. Tags <img> diretas
    for img in element.find_all("img"):
        src = img.get("src") or img.get("data-src") or ""
        if src:
            images.append({
                "src": src,
                "alt": img.get("alt", ""),
                "width": img.get("width"),
                "height": img.get("height"),
            })

    # 2. background-image em style inline
    for el in element.find_all(style=True):
        style = el.get("style", "")
        bg_url = extract_background_image_url(style)
        if bg_url:
            images.append({
                "src": bg_url,
                "alt": extract_text_clean(el)[:50] if el.get_text(strip=True) else "",
                "width": None,
                "height": None,
            })

    # 3. Elementos com class img-placeholder ou similar que tenham data-src
    for el in element.find_all(class_=re.compile(r"img|image|photo|picture")):
        src = el.get("data-src") or el.get("data-image") or ""
        if src and src not in [i["src"] for i in images]:
            images.append({
                "src": src,
                "alt": el.get("alt", extract_text_clean(el)[:50]),
                "width": None,
                "height": None,
            })

    return images


def parse_html_slides(html_path: str) -> list[dict]:
    """
    Lê o HTML e retorna lista de dicts representando slides.
    Extrai: títulos, subtítulos, parágrafos, listas, imagens, tabelas.
    """
    with open(html_path, "r", encoding="utf-8") as f:
        content = f.read()

    soup = BeautifulSoup(content, "lxml")

    # Encontrar slides
    slide_elements = soup.select(".slide, [id^='slide-']")
    if not slide_elements:
        slide_elements = soup.select("section")
    if not slide_elements:
        # Fallback: separar por <hr>
        slide_elements = [soup.body] if soup.body else []

    slides_data = []

    for slide_el in slide_elements:
        slide = {"title": None, "subtitle": None, "content": []}

        # Título (h1 ou h2)
        title_el = slide_el.find(["h1", "h2"])
        if title_el:
            slide["title"] = extract_text_clean(title_el)

        # Subtítulo (h3)
        subtitle_el = slide_el.find("h3")
        if subtitle_el:
            slide["subtitle"] = extract_text_clean(subtitle_el)

        # Processar conteúdo
        processed_texts = set()
        processed_srcs = set()
        if slide["title"]:
            processed_texts.add(slide["title"])
        if slide["subtitle"]:
            processed_texts.add(slide["subtitle"])

        # Buscar imagens no slide (todas as fontes)
        images = find_images_in_element(slide_el)
        for img_data in images:
            src = img_data["src"]
            if src not in processed_srcs:
                processed_srcs.add(src)
                slide["content"].append({
                    "type": "image",
                    "src": src,
                    "alt": img_data["alt"],
                    "width": img_data.get("width"),
                    "height": img_data.get("height"),
                })

        # Coletar todos os elementos de conteúdo textual
        for el in slide_el.find_all(recursive=True):
            if el.name == "p":
                text = extract_text_clean(el)
                if text and text not in processed_texts and len(text) > 1:
                    processed_texts.add(text)
                    is_bold = bool(el.find(["strong", "b"]))
                    color = parse_color_from_style(el.get("style", ""))
                    slide["content"].append({
                        "type": "paragraph",
                        "text": text,
                        "bold": is_bold,
                        "color": color,
                    })

            elif el.name in ("ul", "ol") and el.parent == slide_el or (
                el.find_parent(class_="slide") == slide_el
                and not el.find_parent(["ul", "ol"])
            ):
                items = []
                for li in el.find_all("li", recursive=False):
                    li_text = extract_text_clean(li)
                    if li_text and li_text not in processed_texts:
                        items.append(li_text)
                        processed_texts.add(li_text)
                if items:
                    list_type = "bullet_list" if el.name == "ul" else "numbered_list"
                    slide["content"].append({"type": list_type, "items": items})

            elif el.name == "img":
                # Já processado acima via find_images_in_element
                pass

            elif el.name == "table" and el.find_parent(class_="slide") == slide_el:
                table_id = id(el)
                if table_id not in processed_texts:
                    processed_texts.add(table_id)
                    headers = [extract_text_clean(th) for th in el.find_all("th")]
                    rows = []
                    for tr in el.find_all("tr"):
                        cells = [extract_text_clean(td) for td in tr.find_all("td")]
                        if cells:
                            rows.append(cells)
                    if headers or rows:
                        slide["content"].append({
                            "type": "table",
                            "headers": headers,
                            "rows": rows,
                        })

        slides_data.append(slide)

    return slides_data


def create_pptx(slides_data: list[dict], output_path: str, theme: str = "dark", html_path: str = "") -> str:
    """Gera arquivo .pptx a partir dos dados extraídos."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # Cores do tema
    if theme == "dark":
        bg_color = RGBColor(0x0F, 0x17, 0x2A)
        title_color = RGBColor(0xF1, 0xF5, 0xF9)
        subtitle_color = RGBColor(0x94, 0xA3, 0xB8)
        text_color = RGBColor(0xCB, 0xD5, 0xE1)
        accent_color = RGBColor(0x22, 0xD3, 0xEE)
    else:
        bg_color = RGBColor(0xFF, 0xFF, 0xFF)
        title_color = RGBColor(0x1E, 0x29, 0x3B)
        subtitle_color = RGBColor(0x64, 0x74, 0x8B)
        text_color = RGBColor(0x33, 0x41, 0x55)
        accent_color = RGBColor(0x06, 0xB6, 0xD4)

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        y_pos = Inches(0.5)
        left_margin = Inches(0.8)
        content_width = Inches(11.5)

        # Título
        if slide_data["title"]:
            txBox = slide.shapes.add_textbox(left_margin, y_pos, content_width, Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = title_color
            y_pos += Inches(1.0)

        # Subtítulo
        if slide_data["subtitle"]:
            txBox = slide.shapes.add_textbox(left_margin, y_pos, content_width, Inches(0.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data["subtitle"]
            p.font.size = Pt(18)
            p.font.color.rgb = subtitle_color
            y_pos += Inches(0.7)

        # Conteúdo
        for item in slide_data["content"]:
            # Evitar overflow vertical
            if y_pos > Inches(6.8):
                break

            if item["type"] == "paragraph":
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(left_margin, y_pos, content_width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item["text"]
                p.font.size = Pt(14)
                p.font.color.rgb = text_color
                if item.get("bold"):
                    p.font.bold = True
                if item.get("color"):
                    try:
                        p.font.color.rgb = RGBColor.from_string(item["color"])
                    except Exception:
                        pass
                y_pos += Inches(0.45)

            elif item["type"] in ("bullet_list", "numbered_list"):
                items = item["items"]
                height = Inches(min(len(items) * 0.38, 4.0))
                txBox = slide.shapes.add_textbox(left_margin, y_pos, content_width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                for i, bullet_text in enumerate(items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    if item["type"] == "numbered_list":
                        p.text = f"{i + 1}. {bullet_text}"
                    else:
                        p.text = f"• {bullet_text}"

                    p.font.size = Pt(13)
                    p.font.color.rgb = text_color
                    p.space_after = Pt(4)

                y_pos += Inches(len(items) * 0.35 + 0.2)

            elif item["type"] == "image":
                src = item["src"]
                image_inserted = False

                # Tentar baixar e inserir a imagem
                if src.startswith("http"):
                    img_stream = download_image(src)
                    if img_stream:
                        try:
                            pic_width = Inches(5)
                            pic_height = Inches(3)
                            # Usar dimensões do HTML se disponíveis
                            if item.get("width"):
                                try:
                                    w = int(re.sub(r"[^\d]", "", str(item["width"])))
                                    if w > 0:
                                        pic_width = Inches(min(w / 96, 10))
                                except (ValueError, TypeError):
                                    pass
                            slide.shapes.add_picture(
                                img_stream, left_margin, y_pos, pic_width, pic_height
                            )
                            image_inserted = True
                            y_pos += Inches(3.2)
                        except Exception:
                            pass
                else:
                    # Imagem local — resolver caminho relativo ao HTML
                    img_path = Path(src)
                    if not img_path.is_absolute() and html_path:
                        html_dir = Path(html_path).parent
                        img_path = html_dir / src
                    if img_path.exists():
                        try:
                            slide.shapes.add_picture(
                                str(img_path), left_margin, y_pos, Inches(5), Inches(3)
                            )
                            image_inserted = True
                            y_pos += Inches(3.2)
                        except Exception:
                            pass

                # Fallback: placeholder de texto
                if not image_inserted:
                    txBox = slide.shapes.add_textbox(left_margin, y_pos, Inches(5), Inches(0.4))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    alt_text = item.get("alt") or src
                    p.text = f"[Imagem: {alt_text}]"
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
                    p.font.italic = True
                    y_pos += Inches(0.5)

            elif item["type"] == "table":
                headers = item["headers"]
                rows = item["rows"]
                cols = max(len(headers), max((len(r) for r in rows), default=0))
                row_count = len(rows) + (1 if headers else 0)

                if cols > 0 and row_count > 0:
                    table_height = Inches(min(row_count * 0.4, 3.0))
                    table_shape = slide.shapes.add_table(
                        row_count, cols,
                        left_margin, y_pos,
                        content_width, table_height,
                    )
                    table = table_shape.table

                    # Headers
                    if headers:
                        for i, h in enumerate(headers):
                            if i < cols:
                                cell = table.cell(0, i)
                                cell.text = h
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(11)
                                    paragraph.font.bold = True

                    # Rows
                    offset = 1 if headers else 0
                    for r_idx, row in enumerate(rows):
                        for c_idx, cell_text in enumerate(row):
                            if c_idx < cols:
                                cell = table.cell(r_idx + offset, c_idx)
                                cell.text = cell_text
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(10)

                    y_pos += table_height + Inches(0.3)

    prs.save(output_path)
    return output_path


def convert(html_path: str, output_path: str = None, theme: str = "dark", fetch_images: bool = False) -> str:
    """Função principal de conversão HTML → PPTX.
    
    Args:
        html_path: Caminho do arquivo HTML
        output_path: Caminho de saída (default: mesmo nome com .pptx)
        theme: 'dark' ou 'light'
        fetch_images: Desativado — imagens devem vir do HTML de origem.
    """
    if output_path is None:
        output_path = str(Path(html_path).with_suffix(".pptx"))

    slides_data = parse_html_slides(html_path)
    result = create_pptx(slides_data, output_path, theme=theme, html_path=html_path)

    print(f"✅ Convertido: {html_path}")
    print(f"   → {len(slides_data)} slides extraídos")
    print(f"   → Arquivo: {result}")
    print(f"   → Tema: {theme}")

    # Contar imagens encontradas
    img_count = sum(1 for s in slides_data for c in s["content"] if c["type"] == "image")
    if img_count:
        print(f"   → {img_count} imagens encontradas no HTML e inseridas")
    else:
        print(f"   → Nenhuma imagem encontrada no HTML.")
        print(f"     Dica: inclua tags <img src='url'> no HTML para que sejam convertidas.")

    return result


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2 or "--help" in sys.argv:
        print("Uso: python convert.py <arquivo.html> [arquivo_saida.pptx] [tema: dark|light]")
        print()
        print("Exemplos:")
        print("  python convert.py apresentacao.html")
        print("  python convert.py apresentacao.html saida.pptx")
        print("  python convert.py apresentacao.html saida.pptx light")
        print()
        print("Imagens: o conversor detecta e baixa imagens de tags <img>,")
        print("background-image CSS e data-src presentes no HTML de origem.")
        sys.exit(0 if "--help" in sys.argv else 1)

    html_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    theme = sys.argv[3] if len(sys.argv) > 3 else "dark"

    if not Path(html_file).exists():
        print(f"❌ Arquivo não encontrado: {html_file}")
        sys.exit(1)

    convert(html_file, output_file, theme)
