#!/usr/bin/env python3
"""Convert genapp-politica-publica-apresentacao-light.html to PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
CYAN = RGBColor(0x08, 0x91, 0xB2)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
PINK = RGBColor(0xDB, 0x27, 0x77)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x47, 0x55, 0x69)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_CARD = RGBColor(0xF8, 0xFA, 0xFC)
BORDER_CARD = RGBColor(0xE2, 0xE8, 0xF0)
TEAL_LABEL = RGBColor(0x0E, 0x74, 0x90)


def add_slide_number(slide, num, total):
    """Add slide number at bottom right."""
    txBox = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_card(slide, left, top, width, height, title, body, accent_color=CYAN):
    """Add a card shape with title and body text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = BORDER_CARD
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(14)
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.space_after = Pt(8)
    for line in body.split("\n"):
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRAY
    return shape


def add_flow_box(slide, left, top, text, border_color=CYAN):
    """Add a flow box."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.5), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(4)
    tf.margin_left = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, left, top):
    """Add arrow text between flow boxes."""
    txBox = slide.shapes.add_textbox(left, top, Inches(0.3), Inches(0.7))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(16)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER


def add_title(slide, text, color=DARK):
    """Add a title textbox."""
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    return txBox


def add_subtitle(slide, text, top=Inches(1.2)):
    """Add subtitle text."""
    txBox = slide.shapes.add_textbox(Inches(0.7), top, Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    return txBox


def add_bullet_list(slide, left, top, width, height, items, bullet="→"):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet} {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.space_after = Pt(4)
    return txBox


def add_phase_badge(slide, text, top=Inches(0.3)):
    """Add a phase badge."""
    txBox = slide.shapes.add_textbox(Inches(0.7), top, Inches(8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    return txBox


def build_slide_1(prs):
    """CAPA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(10), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Processo End-to-End"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY

    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11), Inches(2.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Geração de uma Aplicação de IA"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = DARK
    p3 = tf2.add_paragraph()
    p3.text = "no Setor Público"
    p3.font.size = Pt(36)
    p3.font.bold = True
    p3.font.color.rgb = CYAN

    add_subtitle(slide, "Da identificação do problema público à entrega personalizada ao cidadão — o fluxo completo de uma GenApp de Política Pública", Inches(3.5))

    # Tags
    tags = ["GenAI", "Política Pública", "Data Lakehouse", "Guard Rails", "Privacy by Design"]
    tag_colors = [CYAN, PURPLE, GREEN, ORANGE, PINK]
    left = Inches(0.7)
    for i, tag in enumerate(tags):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.8), Inches(1.6), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = tag_colors[i]
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.margin_top = Pt(2)
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(9)
        p.font.color.rgb = tag_colors[i]
        p.alignment = PP_ALIGN.CENTER
        left += Inches(1.8)

    add_slide_number(slide, 1, 17)


def build_slide_2(prs):
    """AGENDA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Agenda")
    col1 = [
        "Visão geral do processo", "Atores envolvidos",
        "Fluxo Inspire — Visão Integrada", "Fase 0 — Compreensão e Design",
        "Fase 1 — Engenharia de Dados", "Fase 2 — Conformidade e Ética",
        "Fase 3 — Fábrica da Solução (GenApp)", "Fase 4 — Entrega e Monitoramento",
    ]
    col2 = [
        "Ciclo de feedback e ajuste", "Arquitetura técnica integrada",
        "Governança e Guard Rails", "Métricas e KPIs",
        "Linha do tempo típica", "Benefícios esperados",
        "Riscos e mitigações", "Próximos passos",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(5.5), col1)
    add_bullet_list(slide, Inches(6.5), Inches(1.5), Inches(5.5), Inches(5.5), col2)
    add_slide_number(slide, 2, 17)


def build_slide_3(prs):
    """VISÃO GERAL"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Visão Geral do Processo")

    # Flow boxes
    boxes = [
        ("Gestor identifica\nproblema", PURPLE), ("Fase 0\nDesign", PURPLE),
        ("Fase 1\nDados", CYAN), ("Fase 2\nÉtica", PINK),
        ("Fase 3\nGenApp", GREEN), ("Fase 4\nEntrega", ORANGE),
        ("Cidadão\natendido", ORANGE),
    ]
    left = Inches(0.3)
    for i, (text, color) in enumerate(boxes):
        add_flow_box(slide, left, Inches(1.5), text, color)
        left += Inches(1.6)
        if i < len(boxes) - 1:
            add_arrow(slide, left, Inches(1.5))
            left += Inches(0.25)

    # Principles
    principles = [
        "Centrado no cidadão (Persona-first)",
        "Dados como ativo estratégico",
        "Ética e privacidade desde o design",
        "IA generativa com Guard Rails",
        "Monitoramento contínuo de impacto",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Princípios do fluxo"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GRAY
    add_bullet_list(slide, Inches(0.7), Inches(3.4), Inches(5.5), Inches(3.5), principles)

    add_card(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(2.5),
             "Resultado Final",
             'Uma "GenApp" — aplicação de IA generativa\nespecializada na política pública — que\nentrega valor personalizado ao cidadão\ne é monitorada em tempo real pelo gestor.')
    add_slide_number(slide, 3, 17)


def build_slide_4(prs):
    """ATORES"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Atores Envolvidos")

    cards_data = [
        ("🏛️ Gestor Público", "Órgão demandante. Identifica o\nproblema, define KPIs, valida\nentregas e monitora resultados.", "Dono do problema e do resultado"),
        ("⚙️ Núcleos CPQD", "7 núcleos especializados que\nexecutam cada fase do processo\n— da prospecção à comunicação.", "Motor técnico e operacional"),
        ("👤 Cidadão (Persona)", "Público-alvo da política pública.\nRecebe a entrega personalizada\nvia canais digitais.", "Beneficiário final e fonte de feedback"),
    ]
    left = Inches(0.5)
    for title, body, role in cards_data:
        add_card(slide, left, Inches(1.5), Inches(3.8), Inches(2.8), title, f"{body}\n\nPapel: {role}")
        left += Inches(4.2)

    # META labels
    metas = ["META3\nNúcleo de IA", "META5\nQualif. Dados", "META1\nInterop.", "META4\nCiber", "META2\nGenAI", "META6\nComun.", "META7\nDestaque"]
    left = Inches(0.3)
    for m in metas:
        add_flow_box(slide, left, Inches(5.2), m, CYAN)
        left += Inches(1.8)
    add_slide_number(slide, 4, 17)


def build_slide_5(prs):
    """FLUXO INSPIRE - IMAGEM"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Fluxo Inspire — Visão Integrada")
    add_subtitle(slide, "Diagrama completo do fluxo de geração da GenApp, desde a identificação do problema até a entrega ao cidadão.")

    img_path = os.path.join(os.path.dirname(__file__), "Fluxo_Inspire.png")
    if os.path.exists(img_path):
        # Center the image
        from PIL import Image
        img = Image.open(img_path)
        img_w, img_h = img.size
        # Max area: 11in wide x 5in tall
        max_w = Inches(11)
        max_h = Inches(5)
        ratio = min(max_w / Emu(int(img_w * 914400 / 96)), max_h / Emu(int(img_h * 914400 / 96)))
        final_w = int(img_w * 914400 / 96 * ratio)
        final_h = int(img_h * 914400 / 96 * ratio)
        left = int((SLIDE_WIDTH - final_w) / 2)
        top = Inches(2.0)
        slide.shapes.add_picture(img_path, left, top, final_w, final_h)
    else:
        # Fallback text if image not found
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "[Imagem Fluxo_Inspire.png não encontrada]"
        p.font.size = Pt(18)
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 5, 17)


def build_slide_6(prs):
    """FASE 0"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_phase_badge(slide, "📋 FASE 0 — Compreensão, Prospecção e Design da Política")
    add_title(slide, "Entendendo o Problema Público")

    items = [
        "Mapeamento de riscos e entendimento profundo do problema",
        "Definição das Personas (público-alvo da política)",
        "Definição de Métricas de Sucesso (KPIs)",
        "Análise de viabilidade técnica e orçamentária",
        "Benchmarking de soluções similares",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(4.0), items)

    add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
             "Saídas desta fase",
             "✅ Escopo aprovado e documentado\n✅ KPIs do projeto definidos e mensuráveis\n✅ Personas mapeadas com necessidades reais\n✅ Riscos identificados e plano de mitigação",
             PURPLE)

    add_card(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(1.8),
             "Por que isso importa",
             "Sem clareza sobre o problema e as personas,\na IA resolve a coisa errada. Esta fase garante\nque a tecnologia serve à política — não o contrário.")
    add_slide_number(slide, 6, 17)


def build_slide_7(prs):
    """FASE 1"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_phase_badge(slide, "🗄️ FASE 1 — Engenharia, Levantamento e Integração de Dados")
    add_title(slide, "Construindo a Base de Dados")

    items_meta5 = [
        "Cruza bases de dados com base nas Personas",
        "Limpeza e normalização no Data Lakehouse",
        "Validação de endereçamento real dos cidadãos",
        "Consolidação de datasets confiáveis",
    ]
    items_meta1 = [
        "Catalogação de metadados da política pública",
        "Disponibilização via APIs REST/GraphQL",
        "Consumo seguro entre sistemas e órgãos",
        "Padrões de interoperabilidade governamental",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "META5 — Qualificação de Dados"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL_LABEL
    add_bullet_list(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(2.0), items_meta5)

    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(5), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "META1 — Interoperabilidade"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = TEAL_LABEL
    add_bullet_list(slide, Inches(0.7), Inches(4.5), Inches(5.5), Inches(2.0), items_meta1)

    add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
             "Saídas desta fase",
             "✅ Datasets consolidados e confiáveis\n✅ Endereçamento real validado\n✅ APIs seguras para consumo (REST/GraphQL)\n✅ Metadados catalogados e documentados",
             CYAN)

    add_card(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(1.8),
             "Arquitetura de Dados",
             "Data Lakehouse centraliza dados brutos e\nprocessados. APIs expõem apenas o necessário,\ncom controle de acesso granular.")
    add_slide_number(slide, 7, 17)


def build_slide_8(prs):
    """FASE 2"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_phase_badge(slide, "🛡️ FASE 2 — Conformidade, Ética e Proteção de Dados")
    add_title(slide, "Garantindo Ética e Segurança")

    items_meta3 = [
        "Aplicação de metodologias de mitigação de vieses",
        "Testes de fairness nos modelos e dados",
        "Documentação de decisões éticas (Model Cards)",
        "Validação de representatividade das Personas",
    ]
    items_meta4 = [
        "Privacy by Design aplicado desde a arquitetura",
        "PETs — Privacy Enhancing Technologies",
        "Anonimização e pseudonimização de dados sensíveis",
        "Conformidade LGPD e Marco Legal da IA",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "META3 — Núcleo de IA (Fase Ética)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL_LABEL
    add_bullet_list(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(2.0), items_meta3)

    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(5), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "META4 — Cibersegurança"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = TEAL_LABEL
    add_bullet_list(slide, Inches(0.7), Inches(4.5), Inches(5.5), Inches(2.0), items_meta4)

    add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
             "Saídas desta fase",
             "✅ Dados do cidadão protegidos (LGPD)\n✅ Infraestrutura validada e segura\n✅ Vieses identificados e mitigados\n✅ Documentação ética completa",
             PINK)

    add_card(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(1.8),
             "Princípio Fundamental",
             "Nenhuma GenApp vai para produção sem\npassar por esta fase. A proteção do cidadão\né pré-requisito, não feature opcional.",
             PINK)
    add_slide_number(slide, 8, 17)


def build_slide_9(prs):
    """FASE 3"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_phase_badge(slide, "🏭 FASE 3 — A Fábrica da Solução e Interpretação Generativa")
    add_title(slide, "Construindo a GenApp")

    items = [
        "Absorve o escopo definido na Fase 0",
        "Consome as APIs seguras das Fases 1 e 2",
        "Orquestra agentes de IA especializados",
        "Calibra Guard Rails no LLM",
        "Configura fluxos de decisão e automação",
        "Testes de qualidade e validação com Personas",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "META2 — Plataforma GenAI (Low/No-code)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL_LABEL
    add_bullet_list(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(3.5), items)

    add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.8),
             "Saída desta fase",
             '✅ "GenApp" da Política Pública\n\nMotor especialista com:\n• Agentes de IA orquestrados\n• Guard Rails calibrados\n• APIs integradas e seguras\n• Pronto para entrega ao cidadão',
             GREEN)

    add_card(slide, Inches(7), Inches(4.9), Inches(5.5), Inches(1.6),
             "Guard Rails",
             "Limites de segurança no LLM que impedem\nrespostas fora do escopo, alucinações,\nvazamento de dados e comportamentos\nnão autorizados.",
             GREEN)
    add_slide_number(slide, 9, 17)


def build_slide_10(prs):
    """FASE 4"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_phase_badge(slide, "📬 FASE 4 — Entrega à Persona e Monitoramento")
    add_title(slide, "Chegando ao Cidadão")

    items_meta6 = [
        "Cruza a GenApp com plataforma CDP",
        "Segmenta por Personas definidas na Fase 0",
        "Gera mensagens proativas e personalizadas",
        "Entrega via Caixa Postal digital (checklist)",
        "Multicanal: app, SMS, e-mail, portal",
    ]
    items_meta7 = [
        "Monitora engajamento na ponta",
        "Acompanha notícias e percepção pública",
        "Mede adesão em tempo real",
        "Gera dashboards e clipping exportável",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "META6 — Comunicação Personalizada"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL_LABEL
    add_bullet_list(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(2.2), items_meta6)

    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(5), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "META7 — DestaqueGovBr"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = TEAL_LABEL
    add_bullet_list(slide, Inches(0.7), Inches(4.8), Inches(5.5), Inches(2.0), items_meta7)

    add_card(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5),
             "Saídas desta fase",
             "✅ Cidadão recebe mensagem proativa\n✅ Checklist na Caixa Postal digital\n✅ Dashboards de monitoramento\n✅ Clipping exportável e insights preditivos",
             ORANGE)

    add_card(slide, Inches(7), Inches(4.6), Inches(5.5), Inches(1.8),
             "CDP — Customer Data Platform",
             "Plataforma que unifica dados do cidadão\nde múltiplas fontes para comunicação\npersonalizada — sem expor dados sensíveis.",
             ORANGE)
    add_slide_number(slide, 10, 17)


def build_slide_11(prs):
    """CICLO DE FEEDBACK"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Ciclo de Feedback e Ajuste Contínuo")

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.5), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "O processo não termina na entrega. O Gestor Público fecha o ciclo comparando métricas reais com os KPIs definidos na Fase 0."
    p.font.size = Pt(13)
    p.font.color.rgb = DARK

    items = [
        "META7 entrega métricas de engajamento e adesão",
        "Gestor compara com KPIs da Fase 0",
        "Identifica gaps e oportunidades de ajuste",
        "Ajusta tom da comunicação (META6)",
        "Refina modelos e Guard Rails (META2/META3)",
        "Ciclo se repete até atingir metas",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(3.5), items)

    add_card(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5),
             "Fluxo Circular",
             "Entrega → Medição → Análise → Ajuste → Entrega\n\n🔄\n\nO gestor garante o sucesso da política\npública através de iteração contínua\nbaseada em dados reais de impacto.")

    add_card(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(1.8),
             "Decisões baseadas em evidência",
             "Dashboards com dados de engajamento real,\nclipping de mídia e insights preditivos\nalimentando cada decisão de ajuste.")
    add_slide_number(slide, 11, 17)


def build_slide_12(prs):
    """ARQUITETURA TÉCNICA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Arquitetura Técnica Integrada")

    stack_items = [
        "Data Lakehouse — Armazenamento unificado",
        "APIs REST/GraphQL — Interoperabilidade",
        "Plataforma GenAI — Low/No-code",
        "LLMs com Guard Rails — Geração controlada",
        "CDP — Personalização de comunicação",
        "PETs — Privacidade computacional",
        "Dashboards BI — Monitoramento em tempo real",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(4.5), stack_items)

    # Layers card
    layers = (
        "▸ Camada de Ingestão\n"
        "  Bases governamentais, registros civis\n\n"
        "▸ Camada de Processamento\n"
        "  Limpeza, cruzamento, anonimização\n\n"
        "▸ Camada de Inteligência\n"
        "  Agentes de IA, LLMs, Guard Rails\n\n"
        "▸ Camada de Entrega\n"
        "  CDP, multicanal, monitoramento"
    )
    add_card(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.0),
             "Camadas da Arquitetura", layers)
    add_slide_number(slide, 12, 17)


def build_slide_13(prs):
    """GOVERNANÇA E GUARD RAILS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Governança e Guard Rails")

    guards = [
        ("Guard Rails de Conteúdo", "LLM não pode gerar informações\nfora do escopo da política."),
        ("Guard Rails de Privacidade", "Dados pessoais nunca são\nexpostos ao modelo."),
        ("Guard Rails de Viés", "Testes de fairness contínuos.\nMonitoramento de disparidade."),
        ("Guard Rails de Segurança", "Proteção contra prompt injection,\ndata leakage e ataques."),
        ("Guard Rails de Escopo", "Agentes operam apenas dentro\ndos limites da Fase 0."),
        ("Guard Rails de Auditoria", "Toda decisão é rastreável.\nLogs completos para accountability."),
    ]

    positions = [
        (Inches(0.5), Inches(1.5)), (Inches(4.5), Inches(1.5)), (Inches(8.5), Inches(1.5)),
        (Inches(0.5), Inches(4.0)), (Inches(4.5), Inches(4.0)), (Inches(8.5), Inches(4.0)),
    ]

    for i, ((title, body), (left, top)) in enumerate(zip(guards, positions)):
        add_card(slide, left, top, Inches(3.8), Inches(2.0), title, body, ORANGE)
    add_slide_number(slide, 13, 17)


def build_slide_14(prs):
    """MÉTRICAS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Métricas e KPIs do Processo")

    # KPIs table as text
    kpis = (
        "Fase 0: Tempo de definição de escopo, clareza de personas\n"
        "Fase 1: Qualidade dos dados, tempo de integração\n"
        "Fase 2: Score de fairness, cobertura LGPD\n"
        "Fase 3: Acurácia do modelo, taxa de Guard Rail triggers\n"
        "Fase 4: Taxa de abertura, engajamento, NPS cidadão"
    )
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(6), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in kpis.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)

    add_card(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.2),
             "Para o Gestor",
             "• % de adesão à política pública\n• Custo por cidadão atendido\n• Tempo médio de resolução\n• ROI da solução vs. processo manual")

    add_card(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(2.2),
             "Para o Cidadão",
             "• Tempo até receber o benefício\n• Satisfação (NPS)\n• Redução de idas presenciais\n• Clareza da comunicação recebida")
    add_slide_number(slide, 14, 17)


def build_slide_15(prs):
    """LINHA DO TEMPO"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Linha do Tempo Típica")

    phases = [
        ("Fase 0", "4-6\nsemanas", "Design + Escopo", PURPLE),
        ("Fase 1", "6-10\nsemanas", "Dados + APIs", CYAN),
        ("Fase 2", "3-5\nsemanas", "Ética + LGPD", PINK),
        ("Fase 3", "8-12\nsemanas", "GenApp + Testes", GREEN),
        ("Fase 4", "4-6\nsemanas", "Entrega + Monitor.", ORANGE),
    ]

    left = Inches(0.5)
    for name, duration, desc, color in phases:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.3), Inches(3.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(12)
        tf.margin_left = Pt(8)
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = duration
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = DARK
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(9)
        p3.font.color.rgb = GRAY
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(8)
        left += Inches(2.5)

    # Total
    add_card(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(1.2),
             "Tempo total estimado: 25 a 39 semanas (6 a 9 meses)",
             "Fases 1 e 2 podem rodar em paralelo, reduzindo para ~5-7 meses em cenário otimizado.")
    add_slide_number(slide, 15, 17)


def build_slide_16(prs):
    """RISCOS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Riscos e Mitigações")

    risks = [
        "Qualidade dos dados de origem (bases desatualizadas)",
        "Resistência institucional à mudança",
        "Dependência de fornecedor único de LLM",
        "Complexidade de integração com sistemas legados",
        "Mudanças regulatórias (Marco Legal da IA)",
        "Expectativas irrealistas de prazo/resultado",
    ]
    mitigations = [
        "META5 valida qualidade antes de prosseguir",
        "Gestão de mudança desde a Fase 0",
        "Arquitetura multi-model (não depender de 1 LLM)",
        "APIs como camada de abstração",
        "Governança adaptativa com revisões trimestrais",
        "KPIs realistas definidos com o gestor na Fase 0",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Riscos Principais"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GRAY
    add_bullet_list(slide, Inches(0.7), Inches(1.8), Inches(5.5), Inches(3.5), risks)

    txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.3), Inches(5), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Estratégias de Mitigação"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = GRAY
    add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(3.5), mitigations)

    add_card(slide, Inches(0.7), Inches(5.5), Inches(11.5), Inches(1.3),
             "Princípio de Segurança",
             "Em caso de dúvida, o sistema para e escala para decisão humana. Nenhum agente de IA toma decisão irreversível sobre a vida do cidadão sem validação.",
             ORANGE)
    add_slide_number(slide, 16, 17)


def build_slide_17(prs):
    """CONCLUSÃO"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Próximos Passos")

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(6), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "O processo de geração de uma GenApp para política pública é estruturado, auditável e centrado no cidadão."
    p.font.size = Pt(14)
    p.font.color.rgb = DARK

    items = [
        "Selecionar política pública piloto (alto impacto, baixo risco)",
        "Ativar META3 para Fase 0 de prospecção",
        "Definir Personas e KPIs com o gestor",
        "Estabelecer governança e comitê de ética",
        "Executar, medir, iterar",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(2.8), Inches(5.5), Inches(3.5), items)

    # Summary card
    summary = (
        "0. Compreensão → Escopo + KPIs\n"
        "1. Dados → Datasets + APIs\n"
        "2. Ética → Proteção + Fairness\n"
        "3. Fábrica → GenApp funcionando\n"
        "4. Entrega → Cidadão atendido\n"
        "🔄 Feedback → Ajuste contínuo"
    )
    add_card(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5),
             "Resumo do Fluxo", summary, TEAL_LABEL)
    add_slide_number(slide, 17, 17)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)
    build_slide_13(prs)
    build_slide_14(prs)
    build_slide_15(prs)
    build_slide_16(prs)
    build_slide_17(prs)

    output_path = os.path.join(os.path.dirname(__file__), "genapp-politica-publica-apresentacao.pptx")
    prs.save(output_path)
    print(f"✅ Apresentação salva em: {output_path}")


if __name__ == "__main__":
    main()
