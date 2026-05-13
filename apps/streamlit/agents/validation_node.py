"""
Validation Node — File verification with retry logic.
Checks that generated files exist, have correct size, and contain expected content.

Model strategy: No LLM needed — deterministic checks.
"""

from pathlib import Path
from bs4 import BeautifulSoup


def validate_html(html_path: str, expected_slides: int) -> dict:
    """
    Validate the generated HTML file.
    Returns dict with is_valid, errors list, and metrics.
    """
    errors = []
    metrics = {}

    path = Path(html_path)

    # Check file exists
    if not path.exists():
        return {"is_valid": False, "errors": ["HTML file not found"], "metrics": {}}

    # Check file size
    size = path.stat().st_size
    metrics["html_size_bytes"] = size

    if size < 500:
        errors.append(f"HTML too small ({size} bytes) — likely empty or broken")
    elif size > 5_000_000:
        errors.append(f"HTML too large ({size} bytes) — possible generation error")

    # Parse and check structure
    try:
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()

        soup = BeautifulSoup(content, "lxml")
        slides = soup.select(".slide, [id^='slide-']")
        metrics["slides_found"] = len(slides)

        if len(slides) == 0:
            errors.append("No slides found in HTML")
        elif len(slides) < expected_slides * 0.5:
            errors.append(
                f"Too few slides: found {len(slides)}, expected ~{expected_slides}"
            )

        # Check navigation exists
        nav = soup.select(".nav, button")
        if not nav:
            errors.append("Navigation buttons not found")

        # Check for empty slides
        empty_count = 0
        for slide in slides:
            text = slide.get_text(strip=True)
            if len(text) < 10:
                empty_count += 1
        metrics["empty_slides"] = empty_count

        if empty_count > len(slides) * 0.3:
            errors.append(f"{empty_count} slides appear empty")

    except Exception as e:
        errors.append(f"HTML parse error: {str(e)}")

    return {
        "is_valid": len(errors) == 0,
        "errors": errors,
        "metrics": metrics,
    }


def validate_pptx(pptx_path: str, expected_slides: int) -> dict:
    """
    Validate the generated PPTX file.
    Returns dict with is_valid, errors list, and metrics.
    """
    errors = []
    metrics = {}

    path = Path(pptx_path)

    # Check file exists
    if not path.exists():
        return {"is_valid": False, "errors": ["PPTX file not found"], "metrics": {}}

    # Check file size
    size = path.stat().st_size
    metrics["pptx_size_bytes"] = size

    if size < 5000:
        errors.append(f"PPTX too small ({size} bytes) — likely corrupted")

    # Check PPTX structure
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        metrics["slides_found"] = slide_count

        if slide_count == 0:
            errors.append("PPTX has no slides")
        elif slide_count < expected_slides * 0.5:
            errors.append(
                f"Too few slides in PPTX: found {slide_count}, expected ~{expected_slides}"
            )

        # Check for text content (not just images)
        text_slides = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    text_slides += 1
                    break
        metrics["slides_with_text"] = text_slides

        if text_slides < slide_count * 0.5:
            errors.append(f"Only {text_slides}/{slide_count} slides have text content")

    except Exception as e:
        errors.append(f"PPTX validation error: {str(e)}")

    return {
        "is_valid": len(errors) == 0,
        "errors": errors,
        "metrics": metrics,
    }


def validation_node(state: dict) -> dict:
    """
    LangGraph node: Validation.
    Verifies generated files and determines if retry is needed.

    Input: html_path, pptx_path, num_slides, formato_saida
    Output: is_valid, error (if any), validation_metrics
    """
    html_path = state.get("html_path", "")
    pptx_path = state.get("pptx_path", "")
    expected_slides = state.get("num_slides", 10)
    formato_saida = state.get("formato_saida", [])
    retry_count = state.get("_retry_count", 0)

    all_errors = []
    all_metrics = {}

    # Validate HTML
    if "HTML (visualização)" in formato_saida or html_path:
        html_result = validate_html(html_path, expected_slides)
        all_errors.extend(html_result["errors"])
        all_metrics["html"] = html_result["metrics"]

    # Validate PPTX
    if "PPTX (editável)" in formato_saida and pptx_path:
        pptx_result = validate_pptx(pptx_path, expected_slides)
        all_errors.extend(pptx_result["errors"])
        all_metrics["pptx"] = pptx_result["metrics"]

    is_valid = len(all_errors) == 0

    # Determine if retry is possible
    can_retry = retry_count < 2 and not is_valid
    error_msg = "; ".join(all_errors) if all_errors else ""

    return {
        **state,
        "is_valid": is_valid,
        "error": error_msg,
        "validation_metrics": all_metrics,
        "_retry_count": retry_count + (1 if not is_valid else 0),
        "_should_retry": can_retry,
    }
