"""
Presentation Generator — Streamlit Interface
Phase 1: Free Presentation Mode
"""

import streamlit as st
from pathlib import Path
import json
import sys

# Ensure config is loaded (dotenv)
sys.path.insert(0, str(Path(__file__).parent))
from config import get_settings

settings = get_settings()

# Page config
st.set_page_config(
    page_title="Gerador de Apresentações",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Custom CSS
st.markdown("""
<style>
    .stApp { background-color: #0f172a; }
    .main .block-container { padding-top: 2rem; max-width: 900px; }
    h1 { color: #f1f5f9 !important; }
    h2, h3 { color: #e2e8f0 !important; }
    .stSelectbox label, .stTextInput label, .stTextArea label, .stSlider label {
        color: #94a3b8 !important;
    }
</style>
""", unsafe_allow_html=True)


# --- Template Registry ---
TEMPLATES_DIR = settings.templates_dir


def load_templates() -> list[dict]:
    """Load available templates from registry."""
    registry_path = TEMPLATES_DIR / "registry.json"
    if registry_path.exists():
        with open(registry_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return []


def get_template_names() -> list[str]:
    """Get list of template display names."""
    templates = load_templates()
    names = ["Padrão (dark theme)"] + [t["name"] for t in templates]
    return names


# --- Sidebar ---
with st.sidebar:
    st.image("https://img.icons8.com/fluency/96/presentation.png", width=64)
    st.title("📊 Apresentações")
    st.markdown("---")

    mode = st.radio(
        "Modo",
        options=["Apresentação Livre", "Proposta de Projeto", "Gerenciar Templates"],
        index=0,
        help="Escolha o tipo de apresentação ou gerencie templates",
    )

    st.markdown("---")
    st.markdown(
        "<small style='color: #64748b;'>Powered by LangGraph + GPT-4o</small>",
        unsafe_allow_html=True,
    )


# --- Main Content ---
st.title("Gerador de Apresentações Profissionais")
st.markdown(
    "<p style='color: #94a3b8; font-size: 1.1rem;'>"
    "Descreva o tema e receba uma apresentação completa com dados reais, "
    "estrutura profissional e exportação editável."
    "</p>",
    unsafe_allow_html=True,
)

st.markdown("---")

if mode == "Apresentação Livre":
    # --- Free Presentation Form ---
    st.subheader("🎯 Briefing da Apresentação")

    col1, col2 = st.columns(2)

    with col1:
        tema = st.text_area(
            "Tema da apresentação",
            placeholder="Ex: O impacto da IA agêntica no setor público brasileiro",
            height=100,
            help="Descreva o tema principal. Quanto mais contexto, melhor o resultado.",
        )

        publico = st.selectbox(
            "Público-alvo",
            options=[
                "Executivos / C-Level",
                "Técnico / Desenvolvedores",
                "Acadêmico / Pesquisadores",
                "Geral / Misto",
                "Investidores / Pitch",
                "Governo / Setor Público",
            ],
            index=0,
        )

        tom = st.selectbox(
            "Tom da apresentação",
            options=[
                "Formal / Corporativo",
                "Técnico / Detalhado",
                "Inspiracional / Motivacional",
                "Casual / Direto",
                "Acadêmico / Científico",
            ],
            index=0,
        )

    with col2:
        num_slides = st.slider(
            "Número de slides",
            min_value=5,
            max_value=25,
            value=12,
            step=1,
            help="Recomendado: 10-15 para apresentações de 20-30 minutos",
        )

        template = st.selectbox(
            "Template visual",
            options=get_template_names(),
            index=0,
            help="Escolha a identidade visual da apresentação",
        )

        formato_saida = st.multiselect(
            "Formatos de exportação",
            options=["HTML (visualização)", "PPTX (editável)"],
            default=["HTML (visualização)", "PPTX (editável)"],
        )

    st.markdown("---")

    # --- Advanced Options (collapsed) ---
    with st.expander("⚙️ Opções avançadas"):
        incluir_referencias = st.checkbox("Incluir slide de referências", value=True)
        incluir_agenda = st.checkbox("Incluir slide de agenda", value=True)
        idioma_pesquisa = st.selectbox(
            "Idioma da pesquisa de dados",
            options=["Português + Inglês", "Apenas Português", "Apenas Inglês"],
            index=0,
        )

    # --- Generate Button ---
    st.markdown("")
    col_btn1, col_btn2, col_btn3 = st.columns([1, 2, 1])

    with col_btn2:
        gerar = st.button(
            "🚀 Gerar Apresentação",
            type="primary",
            use_container_width=True,
            disabled=not tema,
        )

    # --- Generation Flow ---
    if gerar and tema:
        st.markdown("---")

        # Import pipeline nodes
        from agents.research_node import research_node
        from agents.structure_node import structure_node
        from agents.content_node import content_node
        from agents.html_node import html_node
        from agents.pptx_node import pptx_node
        from agents.validation_node import validation_node
        from agents.quality_standards import apply_quality_standards

        # Build initial state
        pipeline_state = {
            "tema": tema,
            "publico": publico,
            "tom": tom,
            "num_slides": num_slides,
            "template": template,
            "formato_saida": formato_saida,
            "incluir_referencias": incluir_referencias,
            "incluir_agenda": incluir_agenda,
            "mode": "free",
            "selected_data": [],
            "research_results": [],
            "research_queries": [],
            "slide_titles": [],
            "slides": [],
            "html_content": "",
            "html_path": "",
            "pptx_path": "",
            "is_valid": False,
            "error": "",
        }

        # Progress container
        progress_container = st.container()

        with progress_container:
            progress_bar = st.progress(0, text="Iniciando...")

            # Step 1: Research
            progress_bar.progress(10, text="🔍 Pesquisando dados em fontes confiáveis...")
            status_research = st.status("Pesquisa de dados", expanded=True)
            with status_research:
                st.write(f"Tema: {tema}")
                st.write(f"Público: {publico}")
                pipeline_state = research_node(pipeline_state)
                n_results = len(pipeline_state.get("research_results", []))
                st.write(f"✅ {n_results} resultados encontrados")
                for r in pipeline_state.get("research_results", [])[:3]:
                    st.caption(f"• [{r.get('source', '')}] {r.get('title', '')[:60]}")
            status_research.update(label=f"Pesquisa — {n_results} fontes", state="complete")

            # Step 2: Structure
            progress_bar.progress(30, text="📐 Estruturando slides...")
            status_structure = st.status("Estruturação", expanded=False)
            with status_structure:
                pipeline_state = structure_node(pipeline_state)
                titles = pipeline_state.get("slide_titles", [])
                st.write(f"✅ {len(titles)} slides estruturados")
                for i, t in enumerate(titles[:5], 1):
                    st.caption(f"{i}. {t}")
                if len(titles) > 5:
                    st.caption(f"... +{len(titles) - 5} slides")
            status_structure.update(label=f"Estruturação — {len(titles)} slides", state="complete")

            # Step 3: Content Generation
            progress_bar.progress(50, text="✍️ Gerando conteúdo...")
            status_content = st.status("Geração de conteúdo", expanded=False)
            with status_content:
                pipeline_state = content_node(pipeline_state)
                slides = pipeline_state.get("slides", [])
                # Apply quality standards
                quality_result = apply_quality_standards(slides)
                pipeline_state["slides"] = quality_result["slides"]
                st.write(f"✅ {len(slides)} slides com conteúdo")
                st.write(f"📊 Quality score: {quality_result['quality_score']}/100")
            status_content.update(label=f"Conteúdo — score {quality_result['quality_score']}/100", state="complete")

            # Step 4: Generate HTML
            progress_bar.progress(70, text="🎨 Gerando apresentação HTML...")
            status_html = st.status("Geração HTML", expanded=False)
            with status_html:
                pipeline_state = html_node(pipeline_state)
                html_path = pipeline_state.get("html_path", "")
                st.write(f"✅ HTML gerado: {Path(html_path).name if html_path else 'erro'}")
            status_html.update(label="HTML gerado", state="complete")

            # Step 5: Convert to PPTX
            if "PPTX (editável)" in formato_saida:
                progress_bar.progress(85, text="📄 Convertendo para PPTX...")
                status_convert = st.status("Conversão PPTX", expanded=False)
                with status_convert:
                    pipeline_state = pptx_node(pipeline_state)
                    pptx_path = pipeline_state.get("pptx_path", "")
                    if pptx_path:
                        size_kb = Path(pptx_path).stat().st_size // 1024
                        st.write(f"✅ PPTX gerado: {Path(pptx_path).name} ({size_kb} KB)")
                    else:
                        st.write(f"⚠️ {pipeline_state.get('error', 'Erro na conversão')}")
                status_convert.update(label="PPTX convertido", state="complete")

            # Step 6: Validation
            progress_bar.progress(95, text="✔️ Validando...")
            pipeline_state = validation_node(pipeline_state)
            progress_bar.progress(100, text="✅ Apresentação gerada!")

        # --- Preview Section ---
        st.markdown("---")
        st.subheader("👁️ Preview da Apresentação")

        html_content = pipeline_state.get("html_content", "")
        if html_content:
            # Tabs: Preview | Slides | Download
            tab_preview, tab_slides, tab_download = st.tabs(["🖥️ Preview", "📋 Slides", "📥 Download"])

            with tab_preview:
                # Render HTML inline
                st.components.v1.html(html_content, height=500, scrolling=True)

            with tab_slides:
                # Show slide content as text
                for slide in pipeline_state.get("slides", []):
                    with st.expander(
                        f"Slide {slide.get('slide_number', '?')}: {slide.get('title', '')}",
                        expanded=False,
                    ):
                        if slide.get("subtitle"):
                            st.caption(slide["subtitle"])
                        for bullet in slide.get("content", []):
                            st.markdown(f"- {bullet}")
                        if slide.get("notes"):
                            st.markdown(f"*Notas: {slide['notes']}*")

            with tab_download:
                col_dl1, col_dl2 = st.columns(2)

                with col_dl1:
                    if html_content:
                        st.download_button(
                            label="⬇️ Download HTML",
                            data=html_content,
                            file_name=Path(html_path).name if html_path else "apresentacao.html",
                            mime="text/html",
                            use_container_width=True,
                        )

                with col_dl2:
                    pptx_path = pipeline_state.get("pptx_path", "")
                    if pptx_path and Path(pptx_path).exists():
                        with open(pptx_path, "rb") as f:
                            pptx_data = f.read()
                        st.download_button(
                            label="⬇️ Download PPTX",
                            data=pptx_data,
                            file_name=Path(pptx_path).name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                        )

                # Validation info
                if pipeline_state.get("is_valid"):
                    st.success("✅ Arquivos validados com sucesso")
                else:
                    st.warning(f"⚠️ {pipeline_state.get('error', 'Validação pendente')}")
        else:
            st.error("❌ Não foi possível gerar a apresentação. Verifique o tema e tente novamente.")

elif mode == "Proposta de Projeto":
    # --- Proposal Mode — Structured Inputs ---
    st.subheader("📋 Proposta de Projeto")
    st.markdown(
        "<p style='color: #94a3b8;'>"
        "Preencha os dados do projeto e receba uma apresentação de proposta "
        "profissional com estrutura fixa: Problema → Solução → Escopo → "
        "Cronograma → Investimento."
        "</p>",
        unsafe_allow_html=True,
    )

    # --- Section 1: Project Identity ---
    st.markdown("#### 🏢 Identificação")
    col_id1, col_id2 = st.columns(2)

    with col_id1:
        nome_projeto = st.text_input(
            "Nome do projeto",
            placeholder="Ex: Plataforma de Gestão Inteligente",
        )
        cliente = st.text_input(
            "Cliente / Organização",
            placeholder="Ex: Ministério da Economia",
        )

    with col_id2:
        responsavel = st.text_input(
            "Responsável pela proposta",
            placeholder="Ex: João Silva — CTO",
        )
        data_proposta = st.date_input("Data da proposta")

    st.markdown("---")

    # --- Section 2: Problem & Solution ---
    st.markdown("#### 🎯 Problema e Solução")

    problema = st.text_area(
        "Problema identificado",
        placeholder=(
            "Descreva o problema que o projeto resolve. "
            "Ex: O órgão processa 5.000 documentos/mês manualmente, "
            "com taxa de erro de 12% e tempo médio de 3 dias por documento."
        ),
        height=100,
    )

    solucao = st.text_area(
        "Solução proposta",
        placeholder=(
            "Descreva a solução em alto nível. "
            "Ex: Sistema de análise documental com IA que automatiza triagem, "
            "classificação e extração de dados, reduzindo tempo para 2 horas."
        ),
        height=100,
    )

    diferenciais = st.text_area(
        "Diferenciais da solução (opcional)",
        placeholder="Ex: IA treinada com dados do setor público, integração com SEI, LGPD compliant",
        height=60,
    )

    st.markdown("---")

    # --- Section 3: Scope ---
    st.markdown("#### 📐 Escopo")

    col_scope1, col_scope2 = st.columns(2)

    with col_scope1:
        escopo_inclui = st.text_area(
            "O que ESTÁ no escopo",
            placeholder=(
                "Liste os entregáveis (um por linha):\n"
                "- Backend com API REST\n"
                "- Frontend web responsivo\n"
                "- Integração com sistema legado X\n"
                "- Treinamento da equipe"
            ),
            height=120,
        )

    with col_scope2:
        escopo_nao_inclui = st.text_area(
            "O que NÃO está no escopo",
            placeholder=(
                "Liste exclusões (um por linha):\n"
                "- Migração de dados históricos\n"
                "- Suporte pós-implantação\n"
                "- Hardware / infraestrutura física"
            ),
            height=120,
        )

    st.markdown("---")

    # --- Section 4: Technology ---
    st.markdown("#### 🛠️ Tecnologia")

    tecnologias = st.text_area(
        "Stack tecnológico",
        placeholder=(
            "Liste as tecnologias principais:\n"
            "- Python + FastAPI (backend)\n"
            "- React + TypeScript (frontend)\n"
            "- LangGraph + GPT-4o (IA)\n"
            "- PostgreSQL (banco de dados)\n"
            "- Docker + AWS (infraestrutura)"
        ),
        height=100,
    )

    st.markdown("---")

    # --- Section 5: Development Phases ---
    st.markdown("#### 📅 Etapas de Desenvolvimento")

    num_etapas = st.number_input(
        "Número de etapas/sprints",
        min_value=2,
        max_value=10,
        value=4,
        step=1,
    )

    etapas = []
    for i in range(int(num_etapas)):
        with st.expander(f"Etapa {i + 1}", expanded=(i == 0)):
            col_et1, col_et2 = st.columns([3, 1])
            with col_et1:
                nome_etapa = st.text_input(
                    "Nome da etapa",
                    placeholder=f"Ex: Sprint {i + 1} — {'MVP' if i == 0 else 'Iteração ' + str(i)}",
                    key=f"etapa_nome_{i}",
                )
                entregas_etapa = st.text_area(
                    "Entregáveis",
                    placeholder="Ex: API de autenticação, tela de login, integração com SSO",
                    height=60,
                    key=f"etapa_entregas_{i}",
                )
            with col_et2:
                duracao_etapa = st.number_input(
                    "Duração (semanas)",
                    min_value=1,
                    max_value=12,
                    value=3,
                    key=f"etapa_duracao_{i}",
                )
            etapas.append({
                "nome": nome_etapa,
                "entregas": entregas_etapa,
                "duracao_semanas": duracao_etapa,
            })

    st.markdown("---")

    # --- Section 6: Investment ---
    st.markdown("#### 💰 Investimento")

    col_inv1, col_inv2 = st.columns(2)

    with col_inv1:
        valor_total = st.text_input(
            "Valor total do projeto",
            placeholder="Ex: R$ 180.000,00",
        )
        forma_pagamento = st.selectbox(
            "Forma de pagamento",
            options=[
                "Por etapa (milestone)",
                "Mensal fixo",
                "50% entrada + 50% entrega",
                "30/30/40 (início/meio/fim)",
                "Personalizado",
            ],
            index=0,
        )

    with col_inv2:
        prazo_total = st.text_input(
            "Prazo total estimado",
            placeholder="Ex: 4 meses",
        )
        validade_proposta = st.text_input(
            "Validade da proposta",
            placeholder="Ex: 15 dias",
            value="15 dias",
        )

    st.markdown("---")

    # --- Section 7: Team (optional) ---
    with st.expander("👥 Equipe (opcional)"):
        equipe = st.text_area(
            "Papéis e profissionais",
            placeholder=(
                "Ex:\n"
                "- Tech Lead — 1 profissional\n"
                "- Desenvolvedor Backend — 2 profissionais\n"
                "- Desenvolvedor Frontend — 1 profissional\n"
                "- Designer UX — 1 profissional (parcial)\n"
                "- Gerente de Projeto — 1 profissional (parcial)"
            ),
            height=100,
        )

    # --- Section 8: Template & Export ---
    st.markdown("---")
    st.markdown("#### 🎨 Apresentação")

    col_tpl1, col_tpl2 = st.columns(2)

    with col_tpl1:
        template_proposta = st.selectbox(
            "Template visual",
            options=get_template_names(),
            index=0,
            key="template_proposta",
        )

    with col_tpl2:
        formato_proposta = st.multiselect(
            "Formatos de exportação",
            options=["HTML (visualização)", "PPTX (editável)"],
            default=["HTML (visualização)", "PPTX (editável)"],
            key="formato_proposta",
        )

    # --- Generate Proposal Button ---
    st.markdown("")
    col_pbtn1, col_pbtn2, col_pbtn3 = st.columns([1, 2, 1])

    with col_pbtn2:
        gerar_proposta = st.button(
            "🚀 Gerar Proposta",
            type="primary",
            use_container_width=True,
            disabled=not (nome_projeto and problema and solucao),
        )

    # --- Proposal Generation Flow ---
    if gerar_proposta and nome_projeto and problema and solucao:
        st.markdown("---")

        # Collect all inputs into a structured dict
        proposal_data = {
            "nome_projeto": nome_projeto,
            "cliente": cliente,
            "responsavel": responsavel,
            "data_proposta": str(data_proposta),
            "problema": problema,
            "solucao": solucao,
            "diferenciais": diferenciais,
            "escopo_inclui": escopo_inclui,
            "escopo_nao_inclui": escopo_nao_inclui,
            "tecnologias": tecnologias,
            "etapas": etapas,
            "valor_total": valor_total,
            "forma_pagamento": forma_pagamento,
            "prazo_total": prazo_total,
            "validade_proposta": validade_proposta,
            "equipe": equipe,
            "template": template_proposta,
        }

        # Progress
        progress_container = st.container()

        with progress_container:
            progress_bar = st.progress(0, text="Iniciando geração da proposta...")

            # Step 1: Structure
            progress_bar.progress(20, text="📐 Estruturando slides da proposta...")
            status_struct = st.status("Estruturação da proposta", expanded=True)
            with status_struct:
                st.write("Framework fixo de proposta:")
                st.write("1. Capa → 2. Problema → 3. Solução → 4. Escopo → "
                         "5. Tecnologia → 6. Etapas → 7. Cronograma → "
                         "8. Equipe → 9. Investimento → 10. Próximos Passos")
                # TODO: Integrate with proposal-generator agent
                st.info("⏳ Integração com agente pendente (Fase 2)")
            status_struct.update(label="Estruturação", state="complete")

            # Step 2: Content generation
            progress_bar.progress(50, text="✍️ Gerando conteúdo persuasivo...")
            status_content = st.status("Redação", expanded=False)
            with status_content:
                st.write("Transformando inputs em texto profissional...")
                st.write("Tom: Formal / Persuasivo")
                # TODO: Integrate with LLM for persuasive writing
                st.info("⏳ Integração com agente pendente (Fase 2)")
            status_content.update(label="Redação", state="complete")

            # Step 3: Generate HTML
            progress_bar.progress(70, text="🎨 Gerando apresentação...")
            status_html = st.status("Geração HTML", expanded=False)
            with status_html:
                st.write(f"Template: {template_proposta}")
                # TODO: Generate HTML with proposal structure
                st.info("⏳ Integração pendente (Fase 2)")
            status_html.update(label="Geração HTML", state="complete")

            # Step 4: Convert
            if "PPTX (editável)" in formato_proposta:
                progress_bar.progress(90, text="📄 Convertendo para PPTX...")
                status_pptx = st.status("Conversão PPTX", expanded=False)
                with status_pptx:
                    st.write("Gerando arquivo editável...")
                    # TODO: Call convert.py
                    st.info("⏳ Integração pendente (Fase 2)")
                status_pptx.update(label="Conversão PPTX", state="complete")

            progress_bar.progress(100, text="✅ Proposta gerada!")

        # --- Results ---
        st.markdown("---")
        st.subheader("📥 Resultado")

        # Show summary of what was captured
        with st.expander("📋 Dados capturados", expanded=False):
            st.json(proposal_data)

        col_pdl1, col_pdl2 = st.columns(2)

        with col_pdl1:
            if "HTML (visualização)" in formato_proposta:
                st.download_button(
                    label="⬇️ Download HTML",
                    data="<!-- Placeholder proposta -->",
                    file_name=f"proposta-{nome_projeto[:30].replace(' ', '-').lower()}.html",
                    mime="text/html",
                    use_container_width=True,
                )

        with col_pdl2:
            if "PPTX (editável)" in formato_proposta:
                st.download_button(
                    label="⬇️ Download PPTX",
                    data=b"placeholder",
                    file_name=f"proposta-{nome_projeto[:30].replace(' ', '-').lower()}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )

        st.success(
            "🎉 Interface de proposta pronta! Na Fase 2, os botões entregarão "
            "arquivos reais gerados pelo agente LangGraph."
        )

elif mode == "Gerenciar Templates":
    # --- Template Management Page ---
    st.subheader("🎨 Gerenciar Templates Corporativos")
    st.markdown(
        "<p style='color: #94a3b8;'>"
        "Faça upload de um arquivo .pptx ou .pdf com a identidade visual da empresa. "
        "O sistema extrai automaticamente cores, fontes e logo para usar nas apresentações."
        "</p>",
        unsafe_allow_html=True,
    )

    st.markdown("---")

    # --- Upload Section ---
    st.markdown("#### ⬆️ Adicionar Novo Template")

    col_up1, col_up2 = st.columns(2)

    with col_up1:
        template_name_input = st.text_input(
            "Nome do template",
            placeholder="Ex: CPQD Institucional",
            help="Nome que aparecerá no dropdown de seleção",
        )

        template_description = st.text_input(
            "Descrição (opcional)",
            placeholder="Ex: Template padrão para apresentações corporativas",
        )

    with col_up2:
        template_footer = st.text_input(
            "Footer (opcional)",
            placeholder="Ex: © 2026 Empresa X — Confidencial",
        )

        uploaded_file = st.file_uploader(
            "Upload do template",
            type=["pptx", "pdf"],
            help="Aceita .pptx (PowerPoint) ou .pdf (será convertido)",
        )

    if uploaded_file and template_name_input:
        if st.button("📥 Registrar Template", type="primary"):
            import tempfile
            import shutil

            with st.spinner("Processando template..."):
                # Save uploaded file temporarily
                suffix = Path(uploaded_file.name).suffix
                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                    tmp.write(uploaded_file.getvalue())
                    tmp_path = tmp.name

                try:
                    if suffix.lower() == ".pdf":
                        # Convert PDF to PPTX first
                        pptx_path = _convert_pdf_to_pptx_template(tmp_path)
                        if not pptx_path:
                            st.error("❌ Não foi possível converter o PDF. Use um arquivo .pptx.")
                            import os
                            os.unlink(tmp_path)
                        else:
                            tmp_path = pptx_path
                            suffix = ".pptx"

                    if suffix.lower() == ".pptx":
                        from template_register import register_template

                        manifest = register_template(
                            pptx_file_path=tmp_path,
                            template_name=template_name_input,
                            description=template_description,
                            footer=template_footer,
                        )

                        st.success(f"✅ Template '{template_name_input}' registrado!")

                        # Show extracted info
                        col_r1, col_r2 = st.columns(2)
                        with col_r1:
                            st.markdown("**Cores extraídas:**")
                            colors = manifest.get("colors", {})
                            for name, hex_val in colors.items():
                                st.markdown(
                                    f"<span style='color:{hex_val}; font-size:1.2rem;'>●</span> "
                                    f"{name}: `{hex_val}`",
                                    unsafe_allow_html=True,
                                )

                        with col_r2:
                            st.markdown("**Fontes:**")
                            fonts = manifest.get("fonts", {})
                            st.write(f"Título: {fonts.get('title', 'N/A')}")
                            st.write(f"Corpo: {fonts.get('body', 'N/A')}")
                            st.write(f"Logo: {'✅ Encontrado' if manifest.get('logo') else '❌ Não encontrado'}")
                            st.write(f"Layouts: {len(manifest.get('layouts', []))}")

                except Exception as e:
                    st.error(f"❌ Erro ao registrar template: {str(e)}")
                finally:
                    import os
                    if os.path.exists(tmp_path):
                        os.unlink(tmp_path)

    st.markdown("---")

    # --- List Existing Templates ---
    st.markdown("#### 📋 Templates Registrados")

    from template_register import list_templates, delete_template

    templates_list = list_templates()

    if not templates_list:
        st.info("Nenhum template corporativo registrado. Use o upload acima para adicionar.")
    else:
        for i, tmpl in enumerate(templates_list):
            with st.expander(f"🎨 {tmpl['name']} ({tmpl['id']})", expanded=False):
                col_t1, col_t2, col_t3 = st.columns([2, 2, 1])

                with col_t1:
                    st.markdown(f"**Descrição:** {tmpl.get('description', 'N/A')}")
                    colors = tmpl.get("colors", {})
                    color_preview = " ".join(
                        f"<span style='color:{v}; font-size:1.5rem;'>●</span>"
                        for v in colors.values()
                    )
                    st.markdown(f"**Paleta:** {color_preview}", unsafe_allow_html=True)

                with col_t2:
                    fonts = tmpl.get("fonts", {})
                    st.markdown(f"**Fonte título:** {fonts.get('title', 'N/A')}")
                    st.markdown(f"**Fonte corpo:** {fonts.get('body', 'N/A')}")

                with col_t3:
                    if st.button("🗑️ Remover", key=f"del_{tmpl['id']}"):
                        delete_template(tmpl["id"])
                        st.rerun()


def _convert_pdf_to_pptx_template(pdf_path: str) -> str | None:
    """
    Convert a PDF to a minimal PPTX template.
    Extracts first page as reference for colors/layout.
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches
        import tempfile

        doc = fitz.open(pdf_path)
        if len(doc) == 0:
            return None

        # Create a PPTX with the first page as image (for color extraction)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Render first page
        page = doc[0]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))

        # Save as temp image
        img_tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        pix.save(img_tmp.name)

        # Add to slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            img_tmp.name, Inches(0), Inches(0),
            prs.slide_width, prs.slide_height,
        )

        # Save PPTX
        pptx_tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(pptx_tmp.name)

        # Cleanup
        import os
        os.unlink(img_tmp.name)
        doc.close()

        return pptx_tmp.name

    except Exception:
        return None
